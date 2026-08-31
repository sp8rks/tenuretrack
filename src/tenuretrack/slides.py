"""The deck (task 8).

Every number on a slide is read from the files in `results/` that the pipeline
wrote. Nothing is recomputed and nothing is retyped, so the deck and the report
cannot disagree: if a slide is wrong, the report is wrong the same way, and one
fix corrects both.

The finished `.pptx` goes through the aggregates-only guardrail before it is
kept. A deck is the artifact most likely to be forwarded to a committee, so it
is the one that most needs to be checked, and it is checked as a file rather
than as an intention.

Layout follows `.claude/skills/deck-builder`. Three rules run through it.

**One grid.** Every slide uses the same margin, the same header block, and the
same footer, so the eye lands in the same place each time and the content is
the only thing that moves. The earlier deck placed each element wherever it
happened to fit, which left a third of most slides empty and gave a reader
nothing to anchor on.

**Say the point in words.** Each slide carries one line under its title stating
what to see. A committee reads a forwarded deck without its presenter, and a
chart that depends on narration arrives silent.

**No default chrome.** No stock table style, no shadows, no gradients. Tables
are a tinted header, hairline rules and alternating rows; numbers sit right of
their labels because that is how a reader compares them down a column.
"""

from __future__ import annotations

import contextlib
import datetime as _dt
import os
import shutil
import subprocess
from collections.abc import Callable, Sequence
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from tenuretrack.branding import logo_path
from tenuretrack.figures import (
    ACCENT,
    BAND,
    INK,
    MUTED,
    PANEL,
    RULE,
    funnel_steps_chart,
    position_chart,
    role_pair_chart,
    trajectory_chart,
    venue_lead_chart,
)
from tenuretrack.guardrail import GuardrailViolation, assert_aggregates_only
from tenuretrack.metrics import METRICS
from tenuretrack.slide_data import (
    SlideData,
    _float,
    load_slide_data,
    subject_slug,
)

__all__ = [
    "SlideData",
    "build_slides",
    "export_pdf",
    "load_slide_data",
    "subject_slug",
]

FIGURES_DIR = "figures"


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.lstrip("#").upper())


INK_RGB = _rgb(INK)
MUTED_RGB = _rgb(MUTED)
ACCENT_RGB = _rgb(ACCENT)
BAND_RGB = _rgb(BAND)
PANEL_RGB = _rgb(PANEL)
RULE_RGB = _rgb(RULE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FACE = "DejaVu Sans"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.62)
CONTENT_W = SLIDE_W - 2 * MARGIN

KICKER_Y = Inches(0.40)
TITLE_Y = Inches(0.62)
RULE_Y = Inches(1.22)
TAKEAWAY_Y = Inches(1.34)
CONTENT_Y = Inches(1.92)
CONTENT_H = Inches(4.62)
FOOTER_RULE_Y = Inches(6.78)
FOOTER_Y = Inches(6.86)

TITLE_SIZE = Pt(25)
TAKEAWAY_SIZE = Pt(13.5)
BODY_SIZE = Pt(12)
SMALL_SIZE = Pt(10)
FOOTER_SIZE = Pt(8.5)

HAIRLINE = Inches(0.012)
TILE_H = Inches(1.12)
NOTE_Y = Inches(6.30)


# ------------------------------------------------------------------ primitives


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _text(
    slide,
    text: str,
    *,
    left,
    top,
    width,
    height,
    size=BODY_SIZE,
    bold=False,
    color=INK_RGB,
    align=PP_ALIGN.LEFT,
    line_spacing: float = 1.25,
    space_after: float = 0.0,
):
    """A text box with no padding, so a stated position is the position."""
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.alignment = align
        para.line_spacing = line_spacing
        para.space_after = Pt(space_after)
        run = para.add_run()
        run.text = line
        run.font.size = size
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FACE
    return box


def _rect(slide, *, left, top, width, height, fill, radius: bool = False):
    """A filled rectangle with the stock shadow turned off.

    python-pptx gives every autoshape PowerPoint's default drop shadow, which
    is the single loudest piece of chrome on an otherwise flat slide.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    if radius:
        # A gentle corner. Some shapes carry no adjustment handle at all.
        with contextlib.suppress(IndexError, ValueError):
            shape.adjustments[0] = 0.06
    return shape


def _rule(slide, *, left, top, width, color=RULE_RGB, thickness=HAIRLINE):
    return _rect(slide, left=left, top=top, width=width, height=thickness, fill=color)


def _picture(slide, path, *, left, top, width, height):
    """Place a figure inside a box, keeping its aspect and centring the slack.

    Handing PowerPoint both a width and a height stretches the image to fit.
    The charts are already drawn at the aspect of the box they belong in, so
    the correction here is usually a few thousandths of an inch, but a figure
    that is silently 3% wide is a figure with the wrong typography.
    """
    from PIL import Image

    with Image.open(path) as img:
        aspect = img.width / img.height
    box_aspect = width / height
    if aspect >= box_aspect:
        draw_w, draw_h = width, int(width / aspect)
    else:
        draw_w, draw_h = int(height * aspect), height
    return slide.shapes.add_picture(
        str(path),
        Emu(int(left + (width - draw_w) / 2)),
        Emu(int(top + (height - draw_h) / 2)),
        width=Emu(int(draw_w)),
        height=Emu(int(draw_h)),
    )


def _cell_borders(cell, color: RGBColor, *, top=False, bottom=False) -> None:
    """Hairlines where they carry meaning, and nowhere else.

    python-pptx has no border API, so the line elements go on by hand. Every
    edge is written, including the ones being switched off: a table that
    inherits its left and right rules from the stock style is a grid, and a
    grid is what makes a table look like a spreadsheet rather than a result.
    """
    properties = cell._tc.get_or_add_tcPr()
    for tag, wanted in (
        ("a:lnL", False), ("a:lnR", False), ("a:lnT", top), ("a:lnB", bottom),
    ):
        for existing in properties.findall(qn(tag)):
            properties.remove(existing)
        line = properties.makeelement(qn(tag), {"w": "9525", "cap": "flat"})
        if wanted:
            fill = line.makeelement(qn("a:solidFill"), {})
            value = fill.makeelement(qn("a:srgbClr"), {"val": str(color)})
            fill.append(value)
            line.append(fill)
        else:
            line.append(line.makeelement(qn("a:noFill"), {}))
        properties.append(line)


def _table(
    slide,
    rows: Sequence[Sequence[str]],
    *,
    left,
    top,
    width,
    row_height,
    widths: Sequence[float] | None = None,
    size=SMALL_SIZE,
):
    """A table with a tinted header, alternating rows and no vertical rules."""
    shape = slide.shapes.add_table(
        len(rows), len(rows[0]), left, top, width, row_height * len(rows)
    )
    table = shape.table
    table.first_row = False
    table.horz_banding = False
    if widths:
        total = sum(widths)
        for column, share in zip(table.columns, widths, strict=True):
            column.width = Emu(int(width * share / total))
    for index in range(len(rows)):
        table.rows[index].height = row_height

    for r, row in enumerate(rows):
        header = r == 0
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                PANEL_RGB if header else (WHITE if r % 2 else _rgb("#fbfcfd"))
            )
            cell.margin_left = cell.margin_right = Inches(0.14)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            _cell_borders(cell, RULE_RGB, top=header, bottom=True)

            frame = cell.text_frame
            frame.word_wrap = True
            for i, line in enumerate(str(value).split("\n")):
                para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
                para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
                para.line_spacing = 1.05
                run = para.add_run()
                run.text = line
                # The second line of a cell is the interval under its number,
                # which is context rather than a value and is set as such.
                secondary = i > 0
                run.font.size = Pt(size.pt - 1.5) if secondary else size
                run.font.bold = header or (c > 0 and not secondary and not header)
                run.font.color.rgb = MUTED_RGB if secondary else INK_RGB
                run.font.name = FACE
    return table


def _tiles(slide, tiles: Sequence[tuple[str, str]], *, top, height=TILE_H):
    """A row of headline numbers, each with the word for what it counts."""
    gap = Inches(0.20)
    slot = int((CONTENT_W - gap * (len(tiles) - 1)) / len(tiles))
    for i, (number, caption) in enumerate(tiles):
        left = MARGIN + i * (slot + gap)
        _rect(slide, left=left, top=top, width=Emu(slot), height=height,
              fill=PANEL_RGB, radius=True)
        _text(
            slide, number, left=left + Inches(0.22), top=top + Inches(0.14),
            width=Emu(slot) - Inches(0.36), height=Inches(0.5),
            size=Pt(23) if len(number) <= 7 else Pt(17), bold=True,
        )
        _text(
            slide, caption, left=left + Inches(0.22), top=top + Inches(0.66),
            width=Emu(slot) - Inches(0.36), height=Inches(0.42),
            size=Pt(9.5), color=MUTED_RGB, line_spacing=1.15,
        )


# ---------------------------------------------------------------- slide frame


def _header(slide, kicker: str, title: str, takeaway: str = "") -> None:
    """The same three lines in the same place on every content slide."""
    _text(slide, kicker.upper(), left=MARGIN, top=KICKER_Y, width=CONTENT_W,
          height=Inches(0.22), size=Pt(9), bold=True, color=ACCENT_RGB)
    _text(slide, title, left=MARGIN, top=TITLE_Y, width=CONTENT_W,
          height=Inches(0.55), size=TITLE_SIZE, bold=True)
    _rule(slide, left=MARGIN, top=RULE_Y, width=Inches(0.86), color=ACCENT_RGB,
          thickness=Inches(0.035))
    if takeaway:
        _text(slide, takeaway, left=MARGIN, top=TAKEAWAY_Y, width=Inches(11.4),
              height=Inches(0.5), size=TAKEAWAY_SIZE, color=MUTED_RGB)


def _footer(slide, data: SlideData, number: int | None) -> None:
    """Provenance and a page number, so a forwarded slide can be traced back."""
    _rule(slide, left=MARGIN, top=FOOTER_RULE_Y, width=CONTENT_W)
    _text(
        slide,
        f"{data.config.subject.name}  |  {data.config.subfield.label}  |  "
        f"cohort of {data.cohort_size:,} at career year {data.horizon}  |  "
        "aggregates only, and a description rather than a standard",
        left=MARGIN, top=FOOTER_Y, width=Inches(11.4), height=Inches(0.3),
        size=FOOTER_SIZE, color=MUTED_RGB,
    )
    if number is not None:
        _text(slide, str(number), left=SLIDE_W - MARGIN - Inches(0.5),
              top=FOOTER_Y, width=Inches(0.5), height=Inches(0.3),
              size=FOOTER_SIZE, color=MUTED_RGB, align=PP_ALIGN.RIGHT)


def _note(slide, text: str, *, top=NOTE_Y) -> None:
    """The one caveat that belongs on this slide, above the footer rule."""
    _text(slide, text, left=MARGIN, top=top, width=Inches(11.6),
          height=Inches(0.42), size=Pt(9.5), color=MUTED_RGB, line_spacing=1.2)


def _speaker(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


# --------------------------------------------------------------------- slides


SHORT_LABELS = {
    "pubs": "Journal articles",
    "led": "Articles led",
    "lead_share": "Share led",
    "citations": "Citations",
    "h_index": "h-index",
    "venue_impact_median": "Median venue impact",
    "top_quartile_share": "Share in top-quartile venues",
}
"""Chart-sized versions of the table labels, for axes with no room for them."""


def _ordered_subject(data: SlideData) -> list[dict]:
    order = [metric.key for metric in METRICS]
    return sorted(
        data.subject,
        key=lambda row: order.index(row["metric"]) if row["metric"] in order else 99,
    )


def _title_slide(prs, data: SlideData, today: _dt.date) -> None:
    slide = _blank(prs)
    subject = data.config.subject
    _rect(slide, left=0, top=0, width=Inches(0.16), height=SLIDE_H, fill=ACCENT_RGB)

    # Top right, and only on this slide: a mark repeated on all eight would be
    # furniture rather than attribution. Sized to finish above 1.66in, which is
    # where the subject's name starts, so a long name cannot run into it.
    logo = logo_path()
    if logo is not None:
        _picture(slide, logo, left=SLIDE_W - MARGIN - Inches(1.55),
                 top=Inches(0.38), width=Inches(1.55), height=Inches(1.18))

    _text(slide, "PUBLICATION NORMS THROUGH THE TENURE CLOCK", left=MARGIN,
          top=Inches(1.28), width=CONTENT_W, height=Inches(0.24), size=Pt(10),
          bold=True, color=ACCENT_RGB)
    _text(slide, subject.name, left=MARGIN, top=Inches(1.66), width=CONTENT_W,
          height=Inches(0.9), size=Pt(40), bold=True)
    _text(slide, f"{data.config.subfield.label}  |  {subject.institution_name}",
          left=MARGIN, top=Inches(2.58), width=CONTENT_W, height=Inches(0.4),
          size=Pt(16), color=MUTED_RGB)
    _text(
        slide,
        f"Career year {data.career_year}, compared against the cohort at "
        f"career year {data.horizon}",
        left=MARGIN, top=Inches(2.96), width=CONTENT_W, height=Inches(0.32),
        size=Pt(13), color=INK_RGB,
    )
    _rule(slide, left=MARGIN, top=Inches(3.40), width=Inches(1.6),
          color=ACCENT_RGB, thickness=Inches(0.04))

    # The same four figures as the report's cover, in the same order. A deck
    # and a report that lead with different numbers read as different studies.
    _tiles(
        slide,
        [
            (f"{data.cohort_size:,}", "people in the cohort"),
            (str(data.horizon), "career year both sides are measured"),
            (f"{data.config.cohort.start_window[0]} to "
             f"{data.config.cohort.start_window[1]}", "when the cohort started"),
            (str(len(data.config.subfield.topics)),
             "OpenAlex topics defining the subfield"),
        ],
        top=Inches(3.76),
    )

    _text(
        slide,
        "These numbers describe what a group of people did. They are not a "
        "standard, nobody in the cohort agreed to be measured, and no part of "
        "this says what any one career should look like.",
        left=MARGIN, top=Inches(5.22), width=Inches(11.0), height=Inches(0.6),
        size=Pt(13), color=INK_RGB,
    )
    _rule(slide, left=MARGIN, top=FOOTER_RULE_Y, width=CONTENT_W)
    _text(
        slide,
        f"Generated {today.isoformat()} with tenuretrack from OpenAlex data "
        "(Priem, Piwowar and Orr, 2022), CC0.  |  Aggregates only: no cohort "
        "member is named anywhere in this deck.",
        left=MARGIN, top=FOOTER_Y, width=Inches(12.0), height=Inches(0.3),
        size=FOOTER_SIZE, color=MUTED_RGB,
    )


def _subject_slide(prs, data: SlideData, figures: Path, number: int) -> None:
    rows = _ordered_subject(data)
    if not rows:
        return
    slide = _blank(prs)
    short = (data.config.subject.name.split() or ["this record"])[0]
    _header(
        slide, "1. Against the cohort",
        f"Where this record sits at year {data.horizon}",
        "The shaded band is the middle half of the cohort. Both sides are "
        f"counted through career year {data.horizon}.",
    )
    chart = position_chart(
        [
            (
                SHORT_LABELS.get(row["metric"], row.get("label", row["metric"])),
                _float(row.get("value")),
                _float(row.get("cohort_p25")),
                _float(row.get("cohort_p50")),
                _float(row.get("cohort_p75")),
                row.get("compared") == "yes",
                row.get("position", ""),
            )
            for row in rows
        ],
        figures / "deck-position.png",
        subject_name=data.config.subject.name,
    )
    _picture(slide, chart, left=MARGIN, top=CONTENT_Y - Inches(0.06),
             width=CONTENT_W, height=Inches(4.32))
    _note(
        slide,
        "A value past either end of a band is drawn just outside it rather than "
        "to scale. Citations carry a count and no position: the cohort's papers "
        f"here are eight to eighteen years old and {short}'s are at most "
        f"{data.horizon}, so the two counts would measure the calendar.",
    )
    _footer(slide, data, number)
    _speaker(
        slide,
        "Subject and cohort are measured at the same career year. A position is "
        "a location in a distribution, not a judgement about a career.",
    )


def _trajectory_slide(prs, data: SlideData, figures: Path, number: int) -> bool:
    by_metric: dict[str, dict[int, dict]] = {}
    for row in data.benchmarks:
        year = _float(row.get("career_year"))
        if year is not None:
            by_metric.setdefault(row["metric"], {})[int(year)] = row
    wanted = ("pubs", "led", "h_index", "lead_share", "top_quartile_share",
              "citations")
    shares = {m.key for m in METRICS if m.is_share}
    labels = {m.key: m.label for m in METRICS}
    short_labels = {
        "pubs": "Journal articles", "led": "Articles led",
        "h_index": "h-index", "lead_share": "Share led",
        "top_quartile_share": "Top-quartile venue share", "citations": "Citations",
    }
    subject_values = {
        row["metric"]: _float(row.get("value"))
        for row in data.subject
        if row.get("compared") == "yes"
    }
    series = []
    for metric in wanted:
        years = sorted(by_metric.get(metric, {}))
        if len(years) < 2:
            continue
        rows = by_metric[metric]
        series.append((
            short_labels.get(metric, labels.get(metric, metric)),
            years,
            [_float(rows[y].get("p25")) or 0.0 for y in years],
            [_float(rows[y].get("p50")) or 0.0 for y in years],
            [_float(rows[y].get("p75")) or 0.0 for y in years],
            subject_values.get(metric),
            metric in shares,
        ))
    if not series:
        return False

    slide = _blank(prs)
    short = (data.config.subject.name.split() or ["this record"])[0]
    _header(
        slide, "2. Across the clock", "The same cohort, year by year",
        "Each panel counts from the start of the appointment through that "
        "career year, so a record earlier on the clock has a year to read "
        "against.",
    )
    chart = trajectory_chart(
        series, short, data.horizon, figures / "deck-trajectory.png"
    )
    _picture(slide, chart, left=MARGIN, top=CONTENT_Y - Inches(0.04),
             width=CONTENT_W, height=Inches(4.26))
    _note(
        slide,
        "Counts are cumulative, so a flat stretch is a year that added little. "
        "Citations are counted as they stand today for the papers published by "
        "that year, which is why the early years look larger here than they did "
        "at the time.",
    )
    _footer(slide, data, number)
    return True


def _norms_slide(prs, data: SlideData, number: int) -> None:
    at_horizon = {
        row["metric"]: row
        for row in data.benchmarks
        if _float(row.get("career_year")) == data.horizon
    }
    if not at_horizon:
        return
    slide = _blank(prs)
    _header(
        slide, f"3. The cohort at year {data.horizon}",
        "The numbers those bands come from",
        f"What {data.cohort_size:,} people had published by the end of career "
        f"year {data.horizon}, as quartiles across people.",
    )

    rows = [[f"Measured through year {data.horizon}", "Lower quarter", "Median",
             "Upper quarter"]]
    for metric in METRICS:
        row = at_horizon.get(metric.key)
        if row is None:
            continue
        rows.append([metric.label] + [_quartile_cell(row, q)
                                      for q in ("p25", "p50", "p75")])

    table_w = Inches(8.3)
    _table(slide, rows, left=MARGIN, top=CONTENT_Y, width=table_w,
           row_height=Inches(0.52), widths=[3.0, 1.5, 1.5, 1.5])

    panel_left = MARGIN + table_w + Inches(0.42)
    panel_w = SLIDE_W - MARGIN - panel_left
    _rect(slide, left=panel_left, top=CONTENT_Y, width=panel_w,
          height=Inches(4.16), fill=PANEL_RGB, radius=True)
    _text(slide, "HOW TO READ A ROW", left=panel_left + Inches(0.28),
          top=CONTENT_Y + Inches(0.28), width=panel_w - Inches(0.56),
          height=Inches(0.24), size=Pt(9), bold=True, color=ACCENT_RGB)
    _text(
        slide,
        "A quarter of the cohort was below the first number.\n\n"
        "Half was below the middle one.\n\n"
        "A quarter was above the last.\n\n"
        "The smaller figures are 95% confidence intervals from a cluster "
        "bootstrap that resamples people, not papers, because one person's "
        "papers are not independent of each other. A wide interval means the "
        "number should be leaned on lightly.",
        left=panel_left + Inches(0.28), top=CONTENT_Y + Inches(0.66),
        width=panel_w - Inches(0.56), height=Inches(3.2), size=Pt(11),
        line_spacing=1.3,
    )
    _note(
        slide,
        "A top-quartile venue is top quartile inside this cohort's own venue "
        "list, not against a global journal ranking, because citation rates "
        "differ enormously between subfields.",
    )
    _footer(slide, data, number)
    _speaker(
        slide,
        "Every figure is a quartile across cohort members. The middle half of "
        "the cohort sat between the first and the last column.",
    )


def _quartile_cell(row: dict, quartile: str) -> str:
    """One quartile with its interval, so a slide carries the report's hedge."""
    value = row.get(quartile) or ""
    if not value:
        return "withheld"
    low, high = row.get(quartile + "_ci_low"), row.get(quartile + "_ci_high")
    if not low or not high:
        return value
    return f"{value}\n{low} to {high}"


def _funnel_slide(prs, data: SlideData, figures: Path, number: int) -> None:
    if not data.funnel:
        return
    slide = _blank(prs)
    _header(
        slide, "4. How it was built", "Who ended up in the cohort",
        "The slide to read first when a cohort looks wrong: a step that removes "
        "almost everybody, or almost nobody, is the one to question.",
    )
    chart = funnel_steps_chart(
        [(label, kept) for label, _rule_text, kept, _dropped in data.funnel],
        figures / "deck-funnel.png",
    )
    _picture(slide, chart, left=MARGIN, top=CONTENT_Y, width=Inches(7.5),
             height=Inches(4.16))

    panel_left = MARGIN + Inches(7.86)
    panel_w = SLIDE_W - MARGIN - panel_left
    _rect(slide, left=panel_left, top=CONTENT_Y, width=panel_w,
          height=Inches(4.16), fill=PANEL_RGB, radius=True)
    _text(slide, "THE SUBFIELD, AS DEFINED HERE", left=panel_left + Inches(0.28),
          top=CONTENT_Y + Inches(0.28), width=panel_w - Inches(0.56),
          height=Inches(0.24), size=Pt(9), bold=True, color=ACCENT_RGB)
    topics = "\n".join(
        f"{topic.name or topic.id}" for topic in data.config.subfield.topics
    ) or "no topics recorded"
    _text(slide, topics, left=panel_left + Inches(0.28),
          top=CONTENT_Y + Inches(0.62), width=panel_w - Inches(0.56),
          height=Inches(1.9), size=Pt(11), line_spacing=1.45)
    _text(
        slide,
        "The cohort is people whose work sits in these topics and whose first "
        "independent appointment could be dated confidently from their bylines. "
        "It is not everyone in the subfield, and the people it drops are not a "
        "random sample of it.",
        left=panel_left + Inches(0.28), top=CONTENT_Y + Inches(2.72),
        width=panel_w - Inches(0.56), height=Inches(1.2), size=Pt(10),
        color=MUTED_RGB, line_spacing=1.3,
    )
    _note(
        slide,
        "Drawn on a log scale so every step stays visible. The exact rule for "
        "each one is in results/funnel.csv and in docs/methods.md.",
    )
    _footer(slide, data, number)


def _venues_slide(prs, data: SlideData, figures: Path, number: int) -> None:
    if not data.venues:
        return
    slide = _blank(prs)
    leads = data.venue_led_share
    _header(
        slide, "5. The venues", "Where this subfield publishes",
        "The journals the cohort used most, so that a top-quartile venue can be "
        "checked against titles rather than taken on trust."
        + ("  The right column is how much of that work the cohort led."
           if leads else ""),
    )
    chart = venue_lead_chart(
        data.venues[:10], leads, figures / "deck-venues.png", horizon=data.horizon
    )
    _picture(slide, chart, left=MARGIN, top=CONTENT_Y - Inches(0.02),
             width=CONTENT_W, height=Inches(4.2))
    _note(
        slide,
        "The figure in brackets is the journal's 2-year mean citedness, a "
        "property of the journal and not of any paper in it. A title near the "
        "top of this list with an impact near zero is usually a conference "
        "abstract series that OpenAlex types as a journal, counted throughout.",
        top=Inches(6.36),
    )
    _footer(slide, data, number)


def _chaperone_slide(prs, data: SlideData, figures: Path, number: int) -> bool:
    if not data.has_chaperone:
        return False
    slide = _blank(prs)
    _header(
        slide, "6. The chaperone effect",
        "Who leads the papers that reach the best venues",
        "Sekara et al. called it the chaperone effect: a route into a selective "
        "venue that runs through a senior co-author. Left is dominated by "
        "whoever wrote the most papers, and on the right every person is their "
        "own control.",
    )
    pooled = [
        (label, rate, data.role_counts.get(label, (0, 0))[1])
        for label, rate in data.role_rates
    ]
    led = data.paired.get("median_led_share")
    middle = data.paired.get("median_middle_share")
    paired = (
        (led, middle, int(data.paired.get("people", 0)))
        if led is not None and middle is not None
        else None
    )
    chart = role_pair_chart(pooled, paired, figures / "deck-roles.png")
    _picture(slide, chart, left=MARGIN, top=CONTENT_Y + Inches(0.06),
             width=CONTENT_W, height=Inches(3.62))

    if data.gap:
        value, low, high = data.gap
        direction = "more often" if value > 0 else "less often"
        interval = (
            f" (95% confidence interval {low:+.1%} to {high:+.1%})"
            if low is not None and high is not None else ""
        )
        _text(
            slide,
            f"Pooled, the cohort's work reached a top-quartile venue "
            f"{abs(value):.1%} {direction} when its members were not leading "
            f"it{interval}. Where the two readings disagree, that disagreement "
            "is the finding.",
            left=MARGIN, top=Inches(5.98), width=Inches(11.6),
            height=Inches(0.4), size=Pt(11.5),
        )
    _note(
        slide,
        "After Sekara et al., PNAS 2018 (doi 10.1073/pnas.1800471115), which "
        "followed authors through time. This is a cross-sectional "
        "approximation: the direction of a difference is informative, its size "
        "should not be read against their figures.",
        top=Inches(6.40),
    )
    _footer(slide, data, number)
    _speaker(
        slide,
        "Led means last author or flagged corresponding. Corresponding flags "
        "are missing for many journal-years, so last position carries most of "
        "the weight.",
    )
    return True


def _caveats_slide(prs, data: SlideData, number: int) -> None:
    slide = _blank(prs)
    _header(
        slide, "Limits", "What this cannot see",
        "Six things to hold in mind before this deck is used for anything.",
    )
    items = [
        ("A publication record is not a person.",
         "Teaching, mentoring, service, funding, software, datasets and public "
         "scholarship do not appear in OpenAlex."),
        ("Citations are shown and never placed.",
         "The cohort's papers here are years older, and citations accumulate "
         "with time."),
        ("Career start is inferred, not looked up.",
         "It comes from publication patterns. Clinical appointments, parental "
         "leave and delayed starts are invisible to it."),
        ("The missing people are not a random sample.",
         "OpenAlex splits some profiles and merges others, so the cohort tilts "
         "toward distinctive names."),
        ("Journal impact says nothing about a paper.",
         "Venue quartiles are computed inside this cohort because citation "
         "cultures differ between subfields."),
        ("This is a description, not an instruction.",
         "A department that reads the median as something everyone must reach "
         "has misread it."),
    ]
    if data.cohort_size and data.cohort_size < 40:
        items.insert(0, (
            f"This cohort is {data.cohort_size} people.",
            "That is small. Read every quartile in it as indicative only.",
        ))

    columns, gap = 2, Inches(0.5)
    column_w = int((CONTENT_W - gap) / columns)
    row_h = Inches(1.34)
    for i, (heading, body) in enumerate(items):
        left = MARGIN + (i % columns) * (column_w + gap)
        top = CONTENT_Y + (i // columns) * row_h
        _rule(slide, left=left, top=top, width=Inches(0.42), color=ACCENT_RGB,
              thickness=Inches(0.026))
        _text(slide, heading, left=left, top=top + Inches(0.16),
              width=Emu(column_w), height=Inches(0.3), size=Pt(12.5), bold=True)
        _text(slide, body, left=left, top=top + Inches(0.52),
              width=Emu(column_w), height=Inches(0.7), size=Pt(10.5),
              color=MUTED_RGB, line_spacing=1.3)

    closing = CONTENT_Y + ((len(items) + 1) // columns) * row_h + Inches(0.10)
    _rect(slide, left=MARGIN, top=closing, width=CONTENT_W, height=Inches(0.62),
          fill=PANEL_RGB, radius=True)
    _text(
        slide,
        "The whole method is in docs/methods.md, and every number here is in "
        "results/ as CSV. Read the funnel before the quartiles.",
        left=MARGIN + Inches(0.28), top=closing + Inches(0.20),
        width=CONTENT_W - Inches(0.56), height=Inches(0.3), size=Pt(11),
    )
    _footer(slide, data, number)


# --------------------------------------------------------------- orchestration


def build_slides(
    data: SlideData,
    results: str | Path,
    *,
    today: _dt.date | None = None,
    on_progress: Callable[[str], None] | None = None,
) -> Path:
    """Write the deck, then prove it carries nothing identifying."""
    results = Path(results)
    figures = results / FIGURES_DIR
    today = today or _dt.date.today()

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    _title_slide(prs, data, today)
    number = 2
    _subject_slide(prs, data, figures, number)
    number += 1
    if _trajectory_slide(prs, data, figures, number):
        number += 1
    _norms_slide(prs, data, number)
    number += 1
    _funnel_slide(prs, data, figures, number)
    number += 1
    _venues_slide(prs, data, figures, number)
    number += 1
    if _chaperone_slide(prs, data, figures, number):
        number += 1
    _caveats_slide(prs, data, number)

    path = results / f"{subject_slug(data.config.subject.name)}.pptx"
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)

    try:
        assert_aggregates_only(path)
    except GuardrailViolation:
        path.unlink(missing_ok=True)
        raise

    if on_progress:
        on_progress(f"Wrote {path.name} ({len(prs.slides)} slides).")
    return path


def export_pdf(
    pptx_path: str | Path, *, on_progress: Callable[[str], None] | None = None
) -> Path | None:
    """Convert to PDF with LibreOffice when it is there, and say so when it is not.

    Never installs anything. A missing converter is a normal state on a fresh
    machine and on Colab, not a failure of the run.
    """
    pptx_path = Path(pptx_path)
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        if on_progress:
            on_progress(
                "PDF export skipped: LibreOffice was not found. Install it and "
                "rerun `tenuretrack slides`, or open the .pptx and export from "
                "there."
            )
        return None

    try:
        subprocess.run(  # noqa: S603
            [
                soffice,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(pptx_path.parent),
                str(pptx_path),
            ],
            check=True,
            capture_output=True,
            timeout=300,
            env={**os.environ, "HOME": str(pptx_path.parent)},
        )
    except (subprocess.SubprocessError, OSError) as exc:
        if on_progress:
            on_progress(f"PDF export failed, leaving the .pptx in place: {exc}")
        return None

    pdf = pptx_path.with_suffix(".pdf")
    if not pdf.exists():
        if on_progress:
            on_progress("PDF export produced no file, leaving the .pptx in place.")
        return None
    if on_progress:
        on_progress(f"Wrote {pdf.name}.")
    return pdf
