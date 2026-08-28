"""The privacy and wording guardrail for everything written under `results/`.

Cohort members are real early-career faculty who never asked to be benchmarked.
Nothing that identifies one of them may leave the local `data/` and `.cache/`
directories. This module is the enforcement point named in CLAUDE.md, and every
writer of a results file calls it before returning.

It also enforces the descriptive-not-prescriptive wording rule: a report states
where a distribution sits, never what someone was supposed to have done.

Violation reports deliberately do not echo the offending name or identifier.
CI captures logs, and a guardrail failure must not itself leak the thing it
caught.
"""

from __future__ import annotations

import csv
import re
import unicodedata
from collections.abc import Iterable, Sequence
from dataclasses import dataclass
from pathlib import Path

__all__ = [
    "GuardrailError",
    "GuardrailViolation",
    "Violation",
    "assert_aggregates_only",
    "assert_directory_aggregates_only",
    "scan_directory",
    "scan_file",
]

AUTHOR_ID_RE = re.compile(r"\bA\d{7,}\b")
ORCID_RE = re.compile(r"\b\d{4}-\d{4}-\d{4}-\d{3}[\dX]\b")

PRESCRIPTIVE_TERMS: tuple[str, ...] = (
    "expected",
    "required",
    "threshold",
    "bar",
    "target",
    "on track",
    "at risk",
    "behind",
    "ahead",
    "quota",
    "minimum",
)
"""Words that turn a description into an instruction. See CLAUDE.md rule 3."""

ALLOWED_PHRASES: tuple[str, ...] = ("minimum cell size",)
"""Phrases that contain a forbidden word but are themselves privacy vocabulary."""

MIN_NAME_LENGTH = 3
"""Shorter strings match too much ordinary text to be usable as a name filter."""

MIN_COHORT_SIZE_FOR_SHAPE_CHECK = 2
"""Below this, a row-count coincidence says nothing about per-person data."""

TEXT_SUFFIXES = frozenset(
    {".md", ".markdown", ".csv", ".tsv", ".txt", ".json", ".yaml", ".yml", ".html", ".svg"}
)
BINARY_SUFFIXES = frozenset({".png", ".jpg", ".jpeg", ".gif", ".pdf", ".xlsx", ".zip"})
PPTX_SUFFIX = ".pptx"


class GuardrailError(RuntimeError):
    """The guardrail could not do its job (unreadable or unscannable input)."""


@dataclass(frozen=True)
class Violation:
    """One guardrail failure. Never carries the offending value itself."""

    path: Path
    rule: str
    detail: str
    line: int | None = None

    def __str__(self) -> str:
        where = f"{self.path}" + (f":{self.line}" if self.line else "")
        return f"{where}: [{self.rule}] {self.detail}"


class GuardrailViolation(Exception):
    """Raised when a results artifact would leak identity or prescribe a number."""

    def __init__(self, violations: Sequence[Violation]) -> None:
        self.violations = list(violations)
        joined = "\n".join(f"  - {v}" for v in self.violations)
        super().__init__(
            f"aggregates-only guardrail failed ({len(self.violations)} violation(s)):"
            f"\n{joined}"
        )


def assert_aggregates_only(
    path: str | Path,
    cohort_names: Iterable[str] = (),
    cohort_ids: Iterable[str] = (),
    *,
    cohort_size: int | None = None,
) -> None:
    """Raise `GuardrailViolation` unless `path` is safe to commit and share."""
    violations = scan_file(
        path, cohort_names, cohort_ids, cohort_size=cohort_size
    )
    if violations:
        raise GuardrailViolation(violations)


def assert_directory_aggregates_only(
    directory: str | Path,
    cohort_names: Iterable[str] = (),
    cohort_ids: Iterable[str] = (),
    *,
    cohort_size: int | None = None,
) -> None:
    """Raise `GuardrailViolation` unless every file under `directory` is safe."""
    violations = scan_directory(
        directory, cohort_names, cohort_ids, cohort_size=cohort_size
    )
    if violations:
        raise GuardrailViolation(violations)


def scan_directory(
    directory: str | Path,
    cohort_names: Iterable[str] = (),
    cohort_ids: Iterable[str] = (),
    *,
    cohort_size: int | None = None,
) -> list[Violation]:
    """Scan every file under a results directory. Returns violations, sorted."""
    directory = Path(directory)
    if not directory.exists():
        return []
    if not directory.is_dir():
        return scan_file(directory, cohort_names, cohort_ids, cohort_size=cohort_size)
    violations: list[Violation] = []
    for path in sorted(p for p in directory.rglob("*") if p.is_file()):
        violations.extend(
            scan_file(path, cohort_names, cohort_ids, cohort_size=cohort_size)
        )
    return violations


def scan_file(
    path: str | Path,
    cohort_names: Iterable[str] = (),
    cohort_ids: Iterable[str] = (),
    *,
    cohort_size: int | None = None,
) -> list[Violation]:
    """Scan one results file. Returns violations rather than raising."""
    path = Path(path)
    if not path.is_file():
        raise GuardrailError(f"nothing to scan at {path}")

    name_patterns = _name_patterns(cohort_names)
    id_patterns = _id_patterns(cohort_ids)
    violations: list[Violation] = []

    text = _read_text(path)
    if text is not None:
        violations.extend(_scan_text(path, text, name_patterns, id_patterns))

    size = cohort_size
    if size is None:
        size = max(len(name_patterns), len(id_patterns))
    if path.suffix.lower() == ".csv" and size >= MIN_COHORT_SIZE_FOR_SHAPE_CHECK:
        violations.extend(_scan_csv_shape(path, size))

    return violations


# ------------------------------------------------------------------ internals


def _read_text(path: Path) -> str | None:
    """Text of a scannable file, or None when there is no text to scan."""
    suffix = path.suffix.lower()
    if suffix in BINARY_SUFFIXES:
        return None
    if suffix == PPTX_SUFFIX:
        return _pptx_text(path)
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        if suffix in TEXT_SUFFIXES:
            raise GuardrailError(
                f"{path} has a text extension but is not valid UTF-8"
            ) from None
        return None
    except OSError as exc:
        raise GuardrailError(f"could not read {path}: {exc}") from exc


def _pptx_text(path: Path) -> str:
    """Every run of text in a deck, one shape per line."""
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise GuardrailError(
            f"cannot scan {path} without python-pptx installed; a deck that cannot "
            "be scanned must not be committed"
        ) from exc
    lines: list[str] = []
    presentation = Presentation(str(path))
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines.append(shape.text_frame.text.replace("\n", " "))
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    lines.append(" ".join(cell.text for cell in row.cells))
    return "\n".join(lines)


def _scan_text(
    path: Path,
    text: str,
    name_patterns: list[re.Pattern[str]],
    id_patterns: list[re.Pattern[str]],
) -> list[Violation]:
    violations: list[Violation] = []
    for number, line in enumerate(text.splitlines(), start=1):
        if AUTHOR_ID_RE.search(line):
            violations.append(
                Violation(path, "openalex_author_id", "an OpenAlex author ID", number)
            )
        if ORCID_RE.search(line):
            violations.append(Violation(path, "orcid", "an ORCID", number))

        folded = _fold(line)
        for pattern in name_patterns:
            if pattern.search(folded):
                violations.append(
                    Violation(
                        path,
                        "cohort_name",
                        "the name of a cohort member",
                        number,
                    )
                )
                break
        for pattern in id_patterns:
            if pattern.search(line):
                violations.append(
                    Violation(
                        path,
                        "cohort_id",
                        "the OpenAlex ID of a cohort member",
                        number,
                    )
                )
                break

        for term, match in _prescriptive_matches(line):
            violations.append(
                Violation(
                    path,
                    "prescriptive_wording",
                    f"the prescriptive word {term!r} (as {match!r}); reports "
                    "describe distributions, they do not set them",
                    number,
                )
            )
    return violations


def _prescriptive_matches(line: str) -> list[tuple[str, str]]:
    haystack = line.lower()
    for phrase in ALLOWED_PHRASES:
        haystack = haystack.replace(phrase, " " * len(phrase))
    found: list[tuple[str, str]] = []
    for term, pattern in _PRESCRIPTIVE_PATTERNS:
        match = pattern.search(haystack)
        if match:
            found.append((term, match.group(0)))
    return found


def _build_prescriptive_patterns() -> list[tuple[str, re.Pattern[str]]]:
    pluralizable = {"threshold", "bar", "target", "quota"}
    out: list[tuple[str, re.Pattern[str]]] = []
    for term in PRESCRIPTIVE_TERMS:
        body = r"\s+".join(re.escape(word) for word in term.split())
        if term in pluralizable:
            body += "s?"
        out.append((term, re.compile(rf"\b{body}\b")))
    return out


_PRESCRIPTIVE_PATTERNS = _build_prescriptive_patterns()


def _fold(value: str) -> str:
    """Casefold and strip accents so `Muller` matches `Müller`."""
    decomposed = unicodedata.normalize("NFKD", value)
    stripped = "".join(c for c in decomposed if not unicodedata.combining(c))
    return stripped.casefold()


def _name_patterns(names: Iterable[str]) -> list[re.Pattern[str]]:
    patterns: list[re.Pattern[str]] = []
    seen: set[str] = set()
    for name in names:
        folded = _fold(str(name)).strip()
        if len(folded) < MIN_NAME_LENGTH or folded in seen:
            continue
        seen.add(folded)
        body = r"\s+".join(re.escape(part) for part in folded.split())
        patterns.append(re.compile(rf"(?<!\w){body}(?!\w)"))
    return patterns


def _id_patterns(ids: Iterable[str]) -> list[re.Pattern[str]]:
    patterns: list[re.Pattern[str]] = []
    seen: set[str] = set()
    for raw in ids:
        short = str(raw).strip().rstrip("/").rsplit("/", 1)[-1]
        if not short or short in seen:
            continue
        seen.add(short)
        patterns.append(re.compile(rf"\b{re.escape(short)}\b"))
    return patterns


def _scan_csv_shape(path: Path, cohort_size: int) -> list[Violation]:
    """A CSV with one row per cohort member is a per-person table in disguise.

    Stripping the names out does not help: metric row plus institution plus
    start year re-identifies people trivially.
    """
    try:
        with path.open(newline="", encoding="utf-8") as handle:
            rows = [row for row in csv.reader(handle) if any(cell.strip() for cell in row)]
    except (OSError, UnicodeDecodeError) as exc:
        raise GuardrailError(f"could not read {path}: {exc}") from exc
    data_rows = max(0, len(rows) - 1)
    if data_rows == cohort_size:
        return [
            Violation(
                path,
                "per_person_table",
                f"{data_rows} data rows for a cohort of {cohort_size} people, "
                "which is a per-person table",
            )
        ]
    return []
