"""The six-slide deck (TASKS.md task 8).

A deck is the artifact most likely to be forwarded to a committee, so the tests
that matter most are the ones about what it must not contain and about it
reading its numbers rather than recomputing them.
"""

from __future__ import annotations

import csv
import datetime as _dt

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from tenuretrack.config import build_config
from tenuretrack.figures import (
    funnel_steps_chart,
    position_chart,
    role_pair_chart,
    trajectory_chart,
    venue_lead_chart,
)
from tenuretrack.guardrail import PRESCRIPTIVE_TERMS, GuardrailViolation
from tenuretrack.slides import build_slides, export_pdf, load_slide_data, subject_slug

# ------------------------------------------------------------------ fixtures


def write_results(tmp_path, *, cohort=900, with_chaperone=True, subject_name=None):
    """A results directory shaped exactly like the pipeline's own output."""
    results = tmp_path / "results"
    results.mkdir(parents=True, exist_ok=True)

    with (results / "funnel.csv").open("w", encoding="utf-8", newline="") as fh:
        w = csv.writer(fh)
        w.writerow(["step", "label", "rule", "kept", "dropped"])
        w.writerow([1, "candidates", "topics T10001", 80000, 0])
        w.writerow([2, "core topic share", "at least 0.4", 5000, 75000])
        w.writerow([3, "start in window", "2008 to 2018", cohort, 5000 - cohort])

    with (results / "benchmarks.csv").open("w", encoding="utf-8", newline="") as fh:
        w = csv.writer(fh)
        w.writerow(["career_year", "metric", "people", "p25", "p50", "p75",
                    "p25_ci_low", "p25_ci_high", "p50_ci_low", "p50_ci_high",
                    "p75_ci_low", "p75_ci_high"])
        for year in (1, 6):
            for metric in ("pubs", "led", "citations"):
                w.writerow([year, metric, cohort, 5, 10, 20, 4, 6, 9, 11, 19, 21])

    with (results / "subject.csv").open("w", encoding="utf-8", newline="") as fh:
        w = csv.writer(fh)
        w.writerow(["career_year", "compared_at", "metric", "label", "value",
                    "cohort_p25", "cohort_p50", "cohort_p75", "position", "compared"])
        w.writerow([11, 6, "pubs", "Journal articles", 30, 5, 10, 20, "above p75", "yes"])
        w.writerow([11, 6, "citations", "Citations", 900, 5, 10, 20, "not compared", "no"])

    with (results / "venues.csv").open("w", encoding="utf-8", newline="") as fh:
        w = csv.writer(fh)
        w.writerow(["venue", "cohort_papers", "impact", "top_quartile"])
        w.writerow(["Chemistry of Materials", 400, "6.3900", "yes"])
        w.writerow(["Journal of Alloys and Compounds", 300, "6.0900", "no"])

    if with_chaperone:
        with (results / "chaperone.csv").open("w", encoding="utf-8", newline="") as fh:
            w = csv.writer(fh)
            w.writerow(["section", "key", "people", "papers", "value", "low", "high"])
            w.writerow(["pooled_rate", "led", 900, 9000, "0.2540", "", ""])
            w.writerow(["pooled_rate", "middle", 900, 12000, "0.2810", "", ""])
    return results


def data_for(tmp_path, config_dict, **kwargs):
    from tenuretrack.report import load_venues

    results = write_results(tmp_path, **kwargs)
    data = load_slide_data(results, build_config(config_dict))
    data.venues = load_venues(results)
    return data, results


# -------------------------------------------------------------------- naming


@pytest.mark.parametrize(
    ("name", "expected"),
    [
        ("Taylor D. Sparks", "taylor-d-sparks"),
        ("Jane Doe", "jane-doe"),
        ("Doe, Jane", "doe-jane"),
        ("", "subject"),
        ("...", "subject"),
    ],
)
def test_subject_slug(name, expected):
    assert subject_slug(name) == expected


# ------------------------------------------------------------------ reading


def test_the_deck_reads_the_pipelines_own_files(tmp_path, config_dict):
    data, _ = data_for(tmp_path, config_dict)
    assert data.cohort_size == 900
    assert data.career_year == 11
    assert data.horizon == 6
    assert [row["metric"] for row in data.subject] == ["pubs", "citations"]
    assert data.venues[0][0] == "Chemistry of Materials"
    assert data.role_rates[0][1] == pytest.approx(0.254)


def test_a_missing_chaperone_file_is_not_an_error(tmp_path, config_dict):
    data, _ = data_for(tmp_path, config_dict, with_chaperone=False)
    assert data.role_rates == []


def test_an_empty_results_directory_reads_as_empty(tmp_path, config_dict):
    data = load_slide_data(tmp_path / "nothing", build_config(config_dict))
    assert data.subject == []
    assert data.cohort_size == 0


# ------------------------------------------------------------------- figures


def test_the_funnel_chart_is_written(tmp_path):
    path = funnel_steps_chart(
        [("candidates", 8000), ("kept", 900)], tmp_path / "f.png"
    )
    assert path.exists() and path.stat().st_size > 1000


def test_the_position_chart_handles_an_uncompared_row(tmp_path):
    rows = [
        ("Journal articles", 30.0, 5.0, 10.0, 20.0, True),
        ("Citations", 900.0, None, None, None, False),
    ]
    path = position_chart(
        [(label, value, p25, p50, p75, compared, "")
         for label, value, p25, p50, p75, compared in rows],
        tmp_path / "d.png",
    )
    assert path.exists()


def test_charts_cope_with_nothing_to_draw(tmp_path):
    assert venue_lead_chart(
        [("Journal", 10, None, False)], {}, tmp_path / "v.png"
    ).exists()
    assert role_pair_chart([("Led", None, 0)], None, tmp_path / "r.png").exists()
    assert trajectory_chart(
        [("Journal articles", [1, 2], [1, 2], [2, 3], [3, 4], None, False)],
        "Jane", 2, tmp_path / "t.png",
    ).exists()


# --------------------------------------------------------------------- deck


def test_the_deck_covers_every_section(tmp_path, config_dict):
    """Title, position, trajectory, norms, funnel, venues, chaperone, limits."""
    data, results = data_for(tmp_path, config_dict)
    deck = build_slides(data, results, today=_dt.date(2026, 8, 28))
    assert len(Presentation(deck).slides) == 8


def test_a_run_without_the_chaperone_pass_drops_that_slide(tmp_path, config_dict):
    """A slide about an analysis that did not run would be a blank assertion."""
    data, results = data_for(tmp_path, config_dict, with_chaperone=False)
    assert len(Presentation(build_slides(data, results)).slides) == 7


def test_the_deck_is_named_after_the_subject(tmp_path, config_dict):
    data, results = data_for(tmp_path, config_dict)
    assert build_slides(data, results).name == "jane-doe.pptx"


def all_text(deck) -> str:
    out = []
    for slide in Presentation(deck).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    out.extend(cell.text for cell in row.cells)
    return "\n".join(out)


def test_the_deck_reads_as_description_not_instruction(tmp_path, config_dict):
    data, results = data_for(tmp_path, config_dict)
    text = all_text(build_slides(data, results)).lower()
    for term in PRESCRIPTIVE_TERMS:
        assert term not in text, f"{term!r} turns a description into an instruction"


def test_the_deck_carries_the_numbers_from_the_files(tmp_path, config_dict):
    data, results = data_for(tmp_path, config_dict)
    text = all_text(build_slides(data, results))
    assert "900" in text  # cohort size, from funnel.csv
    assert "Career year 11" in text  # the subject's own year, from subject.csv
    assert "career year 6" in text  # the comparison horizon, from the same file
    assert "10" in text  # a median, from benchmarks.csv


def test_every_slide_after_the_title_is_numbered_and_attributed(
    tmp_path, config_dict
):
    """A slide gets forwarded on its own, so each one has to say whose it is."""
    slides = list(Presentation(build_slides(*data_for(tmp_path, config_dict))).slides)
    for number, slide in enumerate(slides[1:], start=2):
        text = "\n".join(
            shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
        )
        assert "Jane Doe" in text
        assert "aggregates only" in text
        assert text.rstrip().endswith(str(number))


def test_tables_carry_no_inherited_grid(tmp_path, config_dict):
    """The stock table style rules every edge, which reads as a spreadsheet.

    Only the header's top and each row's bottom are drawn, so this checks the
    left and right edges were explicitly switched off rather than left to
    whatever theme the deck is opened under.
    """
    from pptx.oxml.ns import qn

    deck = build_slides(*data_for(tmp_path, config_dict))
    tables = [
        shape.table
        for slide in Presentation(deck).slides
        for shape in slide.shapes
        if shape.has_table
    ]
    assert tables, "the deck lost its norms table"
    for table in tables:
        for row in table.rows:
            for cell in row.cells:
                properties = cell._tc.find(qn("a:tcPr"))
                for edge in ("a:lnL", "a:lnR"):
                    line = properties.find(qn(edge))
                    assert line is not None
                    assert line.find(qn("a:noFill")) is not None


def test_the_deck_names_only_the_subject(tmp_path, config_dict):
    data, results = data_for(tmp_path, config_dict)
    text = all_text(build_slides(data, results))
    assert "Jane Doe" in text
    assert "A1000001" not in text


def test_the_deck_says_what_it_cannot_see(tmp_path, config_dict):
    data, results = data_for(tmp_path, config_dict)
    text = all_text(build_slides(data, results))
    assert "Teaching" in text
    assert "parental leave" in text
    assert "distinctive names" in text


def test_a_small_cohort_gets_a_caveat_of_its_own(tmp_path, config_dict):
    data, results = data_for(tmp_path, config_dict, cohort=12)
    assert "indicative only" in all_text(build_slides(data, results))


def test_a_big_cohort_does_not(tmp_path, config_dict):
    data, results = data_for(tmp_path, config_dict, cohort=900)
    assert "indicative only" not in all_text(build_slides(data, results))


def test_the_comparison_slides_carry_speaker_notes(tmp_path, config_dict):
    data, results = data_for(tmp_path, config_dict)
    slides = list(Presentation(build_slides(data, results)).slides)
    notes = [
        s.notes_slide.notes_text_frame.text for s in slides if s.has_notes_slide
    ]
    joined = " ".join(notes)
    assert "quartile across cohort members" in joined
    assert "same career year" in joined


def test_no_em_dashes_anywhere_on_a_slide(tmp_path, config_dict):
    data, results = data_for(tmp_path, config_dict)
    assert "—" not in all_text(build_slides(data, results))


def test_a_deck_that_would_leak_is_deleted_not_saved(tmp_path, config_dict):
    """The guardrail runs on the finished file, and a failure leaves nothing
    behind for somebody to find and send on."""
    config_dict["subfield"]["label"] = "work by A5023888391"
    data, results = data_for(tmp_path, config_dict)
    with pytest.raises(GuardrailViolation):
        build_slides(data, results)
    assert not list(results.glob("*.pptx"))


# ----------------------------------------------------------------------- pdf


def test_pdf_export_says_so_when_libreoffice_is_absent(tmp_path, monkeypatch):
    import tenuretrack.slides as slides_mod

    monkeypatch.setattr(slides_mod.shutil, "which", lambda _name: None)
    said: list[str] = []
    assert export_pdf(tmp_path / "deck.pptx", on_progress=said.append) is None
    assert any("LibreOffice was not found" in line for line in said)
    assert any("rerun" in line for line in said)


def test_pdf_export_never_installs_anything(tmp_path, monkeypatch):
    import tenuretrack.slides as slides_mod

    monkeypatch.setattr(slides_mod.shutil, "which", lambda _name: None)

    def explode(*_args, **_kwargs):
        raise AssertionError("must not run a subprocess when soffice is missing")

    monkeypatch.setattr(slides_mod.subprocess, "run", explode)
    assert export_pdf(tmp_path / "deck.pptx") is None


def test_a_failed_conversion_leaves_the_pptx_alone(tmp_path, monkeypatch):
    import subprocess

    import tenuretrack.slides as slides_mod

    deck = tmp_path / "deck.pptx"
    deck.write_bytes(b"not really a deck")
    monkeypatch.setattr(slides_mod.shutil, "which", lambda _name: "/usr/bin/soffice")

    def fail(*_args, **_kwargs):
        raise subprocess.CalledProcessError(1, "soffice")

    monkeypatch.setattr(slides_mod.subprocess, "run", fail)
    said: list[str] = []
    assert export_pdf(deck, on_progress=said.append) is None
    assert deck.exists()
    assert any("leaving the .pptx in place" in line for line in said)


def accent_pixels(png_path) -> int:
    """How many pixels are painted in the subject's accent colour."""
    import matplotlib.image as mpimg
    import numpy as np

    from tenuretrack.figures import ACCENT

    target = np.array(
        [int(ACCENT[i : i + 2], 16) / 255 for i in (1, 3, 5)], dtype=float
    )
    image = mpimg.imread(png_path)[:, :, :3]
    return int((np.abs(image - target).max(axis=2) < 0.02).sum())


def test_an_uncompared_metric_is_drawn_with_no_marker_at_all(tmp_path):
    """A marker anywhere on this row reads as a position, and one in the middle
    reads as "at the median", which is the one thing this row must not say.

    Checked in pixels rather than in the code, because the rule is about what
    a reader sees. The row still gets a band: the cohort's own spread is real
    and worth showing, it is only this record that has nowhere to sit on it.
    """
    uncompared = position_chart(
        [("Citations", 900.0, 5.0, 10.0, 20.0, False, "")], tmp_path / "none.png"
    )
    compared = position_chart(
        [("Journal articles", 30.0, 5.0, 10.0, 20.0, True, "above p75")],
        tmp_path / "some.png",
    )
    assert accent_pixels(uncompared) == 0, "an uncompared row must draw no marker"
    assert accent_pixels(compared) > 0, "a compared row must draw one"


def test_the_logo_lands_on_the_title_slide_and_nowhere_else(tmp_path, config_dict):
    """Attribution once, not furniture on all eight.

    The charts are pictures too, so this counts pictures per slide rather than
    across the deck: the title slide carries no chart, so its one picture is
    the logo.
    """
    data, results = data_for(tmp_path, config_dict)
    slides = Presentation(build_slides(data, results)).slides
    pictures = [
        sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
        for slide in slides
    ]
    assert pictures[0] == 1


def test_a_deck_still_builds_when_the_logo_is_missing(tmp_path, config_dict, monkeypatch):
    """An install stripped of its assets loses a mark, not a deck."""
    monkeypatch.setattr("tenuretrack.slides.logo_path", lambda name=None: None)
    data, results = data_for(tmp_path, config_dict)
    deck = build_slides(data, results)
    slides = Presentation(deck).slides
    assert len(slides) == 8
    assert not [
        shape for shape in slides[0].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
